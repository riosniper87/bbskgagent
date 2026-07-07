"""PPTX -> rasterize each slide to PNG via PowerPoint COM, then hand to VLM.

Uses win32com to open the PPTX in PowerPoint, export each slide as a
high-resolution PNG, and return those images for the VLM to analyze.
Always extracts per-slide text via python-pptx alongside COM images.
Falls back to python-pptx text-only extraction when COM is unavailable.
"""
from __future__ import annotations

import logging
import os
from pathlib import Path

from store_brief.parse.layout_schema import PageRecord
from store_brief.parse.router import ParsedAttachment

log = logging.getLogger(__name__)

MAX_SLIDES = 20
_ppt_app = None


def _get_ppt_app():
    """Lazily create a single PowerPoint COM instance, reused across calls."""
    global _ppt_app
    if _ppt_app is not None:
        try:
            _ = _ppt_app.Presentations.Count
            return _ppt_app
        except Exception:
            _ppt_app = None
    import win32com.client
    app = win32com.client.Dispatch("PowerPoint.Application")
    try:
        app.DisplayAlerts = False
    except Exception:
        pass
    _ppt_app = app
    return app


def _extract_slide_texts(pptx_path: str) -> list[str]:
    from pptx import Presentation

    prs = Presentation(pptx_path)
    texts: list[str] = []
    for slide in prs.slides:
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text.strip())
        texts.append("\n".join(parts))
    return texts


def _rasterize_slides(pptx_path: str, workdir: str, att_id: str) -> list[tuple[str, str]]:
    """Export each slide as PNG via PowerPoint COM. Returns [(path, ref), ...]."""
    app = _get_ppt_app()
    abs_path = str(Path(pptx_path).resolve())
    prs = app.Presentations.Open(abs_path, ReadOnly=True, Untitled=False, WithWindow=False)
    images = []
    try:
        count = min(prs.Slides.Count, MAX_SLIDES)
        for i in range(1, count + 1):
            ref = f"{att_id}#s{i}"
            out = os.path.join(workdir, ref.replace("#", "_") + ".png")
            prs.Slides(i).Export(str(Path(out).resolve()), "PNG", 1280, 720)
            if os.path.isfile(out):
                images.append((out, ref))
    finally:
        prs.Close()
    return images


def _build_pages(att_id: str, slide_texts: list[str], images: list[tuple[str, str]]) -> list[PageRecord]:
    img_by_ref = {ref: path for path, ref in images}
    count = max(len(slide_texts), len(images))
    pages: list[PageRecord] = []
    for i in range(count):
        idx = i + 1
        ref = f"{att_id}#s{idx}"
        pages.append(PageRecord(
            index=idx,
            ref=ref,
            text=slide_texts[i] if i < len(slide_texts) else "",
            image_path=img_by_ref.get(ref),
        ))
    return pages


def _pages_to_text(pages: list[PageRecord]) -> str:
    lines: list[str] = []
    for p in pages:
        lines.append(f"### 슬라이드 {p.index}")
        if p.text.strip():
            lines.append(p.text.strip())
    return "\n".join(lines)


def cleanup_ppt_app():
    """Close the cached PowerPoint COM instance. Call at pipeline end."""
    global _ppt_app
    if _ppt_app is not None:
        try:
            _ppt_app.Quit()
        except Exception:
            pass
        _ppt_app = None


def parse(attachment, workdir: str) -> ParsedAttachment:
    slide_texts = _extract_slide_texts(attachment.path)
    try:
        images = _rasterize_slides(attachment.path, workdir, attachment.id)
        if not images:
            raise RuntimeError("no slides exported")
        pages = _build_pages(attachment.id, slide_texts, images)
        return ParsedAttachment(
            attachment.id,
            text=_pages_to_text(pages),
            image_paths=images,
            pages=pages,
        )
    except Exception as exc:
        log.warning(
            "COM rasterize failed for %s: %s — falling back to text-only",
            attachment.path,
            exc,
        )
        pages = [
            PageRecord(index=i + 1, ref=f"{attachment.id}#s{i + 1}", text=t)
            for i, t in enumerate(slide_texts)
        ]
        return ParsedAttachment(
            attachment.id,
            text=_pages_to_text(pages),
            image_paths=[],
            pages=pages,
        )
