"""Slide-level PPTX normalization (text-only, no COM/VLM)."""
from __future__ import annotations

import logging
import re
from pathlib import Path

from store_brief.ingestion.schema import NormalizedRecord, Provenance
from store_brief.parse.store import StoredParseRecord

log = logging.getLogger(__name__)

MAX_SLIDES = 40
_MIN_BODY_LEN = 12
_BOILERPLATE_RE = re.compile(
    r"^(감사합니다|thank\s*you|목차|contents|표지|cover)$",
    re.I,
)
_BOILERPLATE_PHRASES = (
    "감사합니다",
    "질문 있으시면",
    "문의 바랍니다",
    "본 자료는",
    "무단 배포",
)


def _iter_shapes(shapes):
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    for shp in shapes:
        if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shp.shapes)
        else:
            yield shp


def _shape_text(shape) -> str:
    if shape.has_text_frame and shape.text_frame.text.strip():
        return shape.text_frame.text.strip()
    if getattr(shape, "has_table", False):
        rows: list[str] = []
        for row in shape.table.rows:
            cells = [str(c.text).strip() for c in row.cells if str(c.text).strip()]
            if cells:
                rows.append(" | ".join(cells))
        return "\n".join(rows)
    return ""


# ~5mm rows (EMU): shapes within the same visual row sort left→right instead
# of degenerating into a pure top-then-left order.
_READING_ROW_EMU = 180000


def _slide_reading_order(shapes) -> list:
    return sorted(
        shapes,
        key=lambda s: (
            round((getattr(s, "top", 0) or 0) / _READING_ROW_EMU),
            getattr(s, "left", 0) or 0,
        ),
    )


def _slide_pictures(shapes) -> list:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    return [s for s in shapes if getattr(s, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE]


_FILENAME_DESCR_RE = re.compile(r"^[\w\-. ]+\.(png|jpe?g|gif|bmp|tiff?|emf|wmf)$", re.I)


def _picture_alt_text(shape) -> str:
    """Alt text (descr) of a picture shape, empty if unset or auto-generated."""
    try:
        el = shape._element
        cnvpr = el.nvPicPr.cNvPr
        descr = (cnvpr.get("descr") or "").strip()
    except Exception:
        return ""
    # PowerPoint/python-pptx often auto-fill descr with the image filename —
    # that is not meaningful alt text.
    if _FILENAME_DESCR_RE.match(descr):
        return ""
    return descr


def _is_boilerplate(text: str) -> bool:
    stripped = (text or "").strip()
    if len(stripped) < _MIN_BODY_LEN:
        return True
    first_line = stripped.split("\n", 1)[0].strip()
    if _BOILERPLATE_RE.match(first_line):
        return True
    if len(stripped) < 40 and any(p in stripped for p in _BOILERPLATE_PHRASES):
        return True
    return False


def _slide_title_and_body(parts: list[str]) -> tuple[str, str]:
    if not parts:
        return "", ""
    title = parts[0][:120] if parts else ""
    body = "\n".join(parts[1:]) if len(parts) > 1 else parts[0]
    return title, body


def parse_pptx_path(path: str, post_id: str, attachment_id: str, filename: str) -> list[NormalizedRecord]:
    """Parse PPTX file into slide-level NormalizedRecords."""
    from pptx import Presentation

    prs = Presentation(path)
    records: list[NormalizedRecord] = []
    for i, slide in enumerate(prs.slides):
        if i >= MAX_SLIDES:
            break
        idx = i + 1
        source_ref = f"{attachment_id}#s{idx}"
        shapes = list(_iter_shapes(slide.shapes))
        ordered = _slide_reading_order(shapes)
        parts = [_shape_text(s) for s in ordered]
        parts = [p for p in parts if p]

        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
        if notes:
            parts.append(f"[노트] {notes}")

        combined = "\n".join(parts).strip()

        pictures = _slide_pictures(shapes)
        if not combined and pictures:
            # Image-only slide: emit a flagged record instead of dropping it.
            # Body stays alt-text-only (or empty) so the VLM description
            # pipeline can supply text via source_ref downstream.
            alt_texts = [t for t in (_picture_alt_text(p) for p in pictures) if t]
            records.append(NormalizedRecord(
                post_id=post_id,
                source_type="pptx_slide",
                title=f"슬라이드 {idx} (이미지)",
                body="\n".join(alt_texts),
                attachment_name=filename,
                provenance=Provenance(
                    source_ref=source_ref,
                    locator=f"슬라이드 {idx}",
                    raw={"slide": idx, "images": len(pictures)},
                    extraction="fallback",
                ),
                review_flag="pptx_image_only",
            ))
            continue

        if _is_boilerplate(combined):
            continue

        title, body = _slide_title_and_body(parts)
        if not body:
            body = combined

        review_flag = None
        extraction: str = "deterministic"

        records.append(NormalizedRecord(
            post_id=post_id,
            source_type="pptx_slide",
            title=title or f"슬라이드 {idx}",
            body=body,
            attachment_name=filename,
            provenance=Provenance(
                source_ref=source_ref,
                locator=f"슬라이드 {idx}",
                raw={"slide": idx},
                extraction=extraction,  # type: ignore[arg-type]
            ),
            review_flag=review_flag,
        ))
    return records


def parse_pptx_from_record(record: StoredParseRecord, post_id: str) -> list[NormalizedRecord]:
    """Parse stored PPTX record using file path or embedded page text."""
    path = record.source_path
    if path and Path(path).is_file():
        try:
            return parse_pptx_path(path, post_id, record.attachment_id, record.filename)
        except Exception as exc:
            log.warning("pptx parse failed for %s: %s", record.filename, exc)

    if record.pages:
        records: list[NormalizedRecord] = []
        for page in record.pages:
            ref = page.get("ref") if isinstance(page, dict) else page.ref
            text = (page.get("text") if isinstance(page, dict) else page.text) or ""
            if _is_boilerplate(text):
                continue
            m = re.search(r"#s(\d+)$", ref or "")
            idx = int(m.group(1)) if m else len(records) + 1
            records.append(NormalizedRecord(
                post_id=post_id,
                source_type="pptx_slide",
                title=f"슬라이드 {idx}",
                body=text.strip(),
                attachment_name=record.filename,
                provenance=Provenance(
                    source_ref=ref or f"{record.attachment_id}#s{idx}",
                    locator=f"슬라이드 {idx}",
                    raw={"slide": idx},
                    extraction="deterministic",
                ),
            ))
        return records

    # fallback: split stored text by slide markers
    blocks = re.split(r"(?=###\s*슬라이드\s*\d+)", record.text or "")
    out: list[NormalizedRecord] = []
    for block in blocks:
        block = block.strip()
        if not block:
            continue
        m = re.search(r"###\s*슬라이드\s*(\d+)", block)
        if not m:
            continue
        idx = int(m.group(1))
        body = re.sub(r"###\s*슬라이드\s*\d+\s*", "", block).strip()
        if _is_boilerplate(body):
            continue
        out.append(NormalizedRecord(
            post_id=post_id,
            source_type="pptx_slide",
            title=f"슬라이드 {idx}",
            body=body,
            attachment_name=record.filename,
            provenance=Provenance(
                source_ref=f"{record.attachment_id}#s{idx}",
                locator=f"슬라이드 {idx}",
                raw={"slide": idx},
                extraction="deterministic",
            ),
        ))
    return out
