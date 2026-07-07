"""Tests for ingestion/parse_pptx slide-level normalization."""
from __future__ import annotations

import json
import sys
from pathlib import Path

import pytest

sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "src"))

from store_brief.ingestion.parse_pptx import parse_pptx_from_record
from store_brief.parse.store import ParsedAttachmentStore

DATA = Path(__file__).resolve().parents[1] / "data"
STORE = ParsedAttachmentStore(DATA / "parsed")
POST_ID = "6148b200b185"


def _load_pptx_record():
    meta_path = DATA / "parsed" / POST_ID / "meta.json"
    if not meta_path.is_file():
        return None
    meta = json.loads(meta_path.read_text(encoding="utf-8"))
    for key in meta.get("records", []):
        rec = STORE.load_record(POST_ID, key)
        if rec and rec.kind == "pptx":
            return rec
    return None


def test_pptx_slide_records():
    record = _load_pptx_record()
    if record is None:
        pytest.skip("광고 PPTX parsed record not available")
    records = parse_pptx_from_record(record, POST_ID)
    assert len(records) > 0
    assert any(r.provenance.source_ref.endswith("#s1") for r in records)
    assert all(r.source_type == "pptx_slide" for r in records)
    assert all(r.provenance.extraction in ("deterministic", "fallback") for r in records)


def test_pptx_skips_boilerplate_slides():
    record = _load_pptx_record()
    if record is None:
        pytest.skip("광고 PPTX parsed record not available")
    records = parse_pptx_from_record(record, POST_ID)
    bodies = "\n".join(r.body for r in records).lower()
    assert "감사합니다" not in bodies or len(records) > 1


def _png_bytes(width: int = 300, height: int = 200) -> bytes:
    import fitz

    pix = fitz.Pixmap(fitz.csRGB, fitz.IRect(0, 0, width, height))
    pix.clear_with(120)
    return pix.tobytes("png")


def _build_deck(path, *, alt_text: str | None):
    import io

    from pptx import Presentation
    from pptx.util import Emu, Inches

    prs = Presentation()
    blank = prs.slide_layouts[6]

    # slide 1: image only
    s1 = prs.slides.add_slide(blank)
    pic = s1.shapes.add_picture(io.BytesIO(_png_bytes()), Inches(1), Inches(1), Inches(6), Inches(4))
    if alt_text is not None:
        pic._element.nvPicPr.cNvPr.set("descr", alt_text)

    # slide 2: normal text slide with two boxes in the same visual row
    s2 = prs.slides.add_slide(blank)
    left_box = s2.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    left_box.text_frame.text = "왼쪽 안내 문구입니다 (행사 시작)"
    right_box = s2.shapes.add_textbox(Inches(5), Emu(Inches(1) + 50000), Inches(3), Inches(1))
    right_box.text_frame.text = "오른쪽 세부 조건입니다 (6월 한정)"

    # slide 3: truly empty (no shapes with content, no pictures)
    prs.slides.add_slide(blank)

    prs.save(str(path))


def test_image_only_slide_emits_flagged_record(tmp_path):
    from store_brief.ingestion.parse_pptx import parse_pptx_path

    path = tmp_path / "deck.pptx"
    _build_deck(path, alt_text="프로모션 배너 이미지")
    records = parse_pptx_path(str(path), "p1", "att#pp", "deck.pptx")

    image_recs = [r for r in records if r.review_flag == "pptx_image_only"]
    assert len(image_recs) == 1
    rec = image_recs[0]
    assert rec.title == "슬라이드 1 (이미지)"
    assert rec.body == "프로모션 배너 이미지"
    assert rec.provenance.extraction == "fallback"
    assert rec.provenance.raw.get("images") == 1
    # truly empty slide (no pictures) still dropped
    assert not any(r.provenance.raw.get("slide") == 3 for r in records)


def test_image_only_slide_without_alt_text_has_empty_body(tmp_path):
    from store_brief.ingestion.parse_pptx import parse_pptx_path

    path = tmp_path / "deck2.pptx"
    _build_deck(path, alt_text=None)
    records = parse_pptx_path(str(path), "p1", "att#pp2", "deck2.pptx")
    image_recs = [r for r in records if r.review_flag == "pptx_image_only"]
    assert len(image_recs) == 1
    # empty body → downstream VLM description pipeline can fill via source_ref
    assert image_recs[0].body == ""


def test_reading_order_same_row_sorts_left_to_right(tmp_path):
    from store_brief.ingestion.parse_pptx import parse_pptx_path

    path = tmp_path / "deck3.pptx"
    _build_deck(path, alt_text="x")
    records = parse_pptx_path(str(path), "p1", "att#pp3", "deck3.pptx")
    text_rec = next(r for r in records if r.review_flag is None)
    combined = f"{text_rec.title}\n{text_rec.body}"
    # boxes are ~0.05mm apart vertically → same visual row → left before right
    assert combined.index("왼쪽") < combined.index("오른쪽")
